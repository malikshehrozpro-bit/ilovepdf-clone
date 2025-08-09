import argparse
import os
import json
import fitz  # PyMuPDF
from typing import List
from PIL import Image
import pikepdf
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
import tabula
import pandas as pd

def merge_pdfs(inputs: List[str], output: str):
  doc = fitz.open()
  for p in inputs:
    with fitz.open(p) as src:
      doc.insert_pdf(src)
  doc.save(output)
  doc.close()

def parse_ranges(ranges: str, max_pages: int) -> List[int]:
  pages = set()
  if not ranges:
    return list(range(1, max_pages + 1))
  for token in ranges.split(","):
    token = token.strip()
    if "-" in token:
      a, b = token.split("-", 1)
      try:
        a_i, b_i = int(a), int(b)
        for i in range(a_i, b_i + 1):
          if 1 <= i <= max_pages: pages.add(i)
      except:
        continue
    else:
      try:
        i = int(token)
        if 1 <= i <= max_pages: pages.add(i)
      except:
        continue
  return sorted(pages)

def split_pdf(input_path: str, ranges: str, outdir: str):
  os.makedirs(outdir, exist_ok=True)
  src = fitz.open(input_path)
  pages = parse_ranges(ranges, src.page_count)
  if not pages:
    pages = list(range(1, src.page_count + 1))
  # either individual files per page or grouped by ranges provided
  # Here: export each selected page as separate PDF
  for p in pages:
    doc = fitz.open()
    doc.insert_pdf(src, from_page=p-1, to_page=p-1)
    out = os.path.join(outdir, f"page-{p}.pdf")
    doc.save(out)
    doc.close()
  src.close()

def compress_pdf(input_path: str, output: str):
  # Simple re-save to optimize streams; for real compression use Ghostscript (handled in Node)
  doc = fitz.open(input_path)
  doc.save(output, garbage=4, deflate=True)
  doc.close()

def pdf_to_word(input_path: str, output: str):
  cv = Converter(input_path)
  cv.convert(output, start=0, end=None)
  cv.close()

def pdf_to_pptx(input_path: str, output: str):
  prs = Presentation()
  with fitz.open(input_path) as doc:
    for i, page in enumerate(doc):
      pix = page.get_pixmap(dpi=200, alpha=False)
      img_path = output + f".tmp_{i}.jpg"
      pix.save(img_path)
      blank = prs.slide_layouts[6]
      slide = prs.slides.add_slide(blank)
      # Set slide size to match image roughly (use 10in x 7.5in default, fit image)
      left = Inches(0)
      top = Inches(0)
      slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
      os.remove(img_path)
  prs.save(output)

def jpg_to_pdf(images: List[str], output: str):
  pdf = fitz.open()
  for img_path in images:
    img = Image.open(img_path).convert("RGB")
    # Create a page with image size at 72dpi scale
    rect = fitz.Rect(0, 0, img.width, img.height)
    page = pdf.new_page(width=rect.width, height=rect.height)
    img_bytes = open(img_path, "rb").read()
    pix = fitz.Pixmap(img_bytes)
    page.insert_image(rect, stream=img_bytes)
  pdf.save(output)
  pdf.close()

def pdf_to_jpg(input_path: str, outdir: str):
  os.makedirs(outdir, exist_ok=True)
  with fitz.open(input_path) as doc:
    for i, page in enumerate(doc):
      pix = page.get_pixmap(dpi=200, alpha=False)
      out = os.path.join(outdir, f"page-{i+1}.jpg")
      pix.save(out)

def rotate_pdf(input_path: str, output: str, rotation: int, ranges: str):
  doc = fitz.open(input_path)
  pages = parse_ranges(ranges, doc.page_count) if ranges else list(range(1, doc.page_count + 1))
  for p in pages:
    page = doc.load_page(p - 1)
    page.set_rotation(((page.rotation + rotation) % 360))
  doc.save(output)
  doc.close()

def unlock_pdf(input_path: str, output: str, password: str):
  # Try open with password and save without encryption
  with pikepdf.open(input_path, password=password) as pdf:
    pdf.save(output)

def protect_pdf(input_path: str, output: str, password: str):
  with pikepdf.open(input_path) as pdf:
    pdf.save(output, encryption=pikepdf.Encryption(owner=password, user=password, R=4))

def reorder_pages(input_path: str, output: str, order: List[int]):
  doc = fitz.open(input_path)
  new = fitz.open()
  if not order:
    order = list(range(0, doc.page_count))
  for idx in order:
    if 0 <= idx < doc.page_count:
      new.insert_pdf(doc, from_page=idx, to_page=idx)
  new.save(output)
  new.close()
  doc.close()

def watermark_text(input_path: str, output: str, text: str):
  doc = fitz.open(input_path)
  for page in doc:
    rect = page.rect
    page.insert_textbox(rect, text or "CONFIDENTIAL",
                        fontsize=48, fontname="helv",
                        rotate=45, color=(0.7,0,0), fill_opacity=0.2)
  doc.save(output)
  doc.close()

def watermark_image(input_path: str, output: str, image_path: str):
  doc = fitz.open(input_path)
  for page in doc:
    rect = page.rect
    margin = 36
    box = fitz.Rect(margin, margin, rect.width - margin, rect.height - margin)
    page.insert_image(box, filename=image_path, keep_proportion=True, overlay=True, opacity=0.2)
  doc.save(output)
  doc.close()

def pdf_to_excel(input_path: str, output: str):
  # Requires Java. Extract tables from all pages and write each as separate sheet
  dfs = tabula.read_pdf(input_path, pages="all", multiple_tables=True, lattice=False, stream=True)
  if not dfs or len(dfs) == 0:
    # Write empty workbook
    with pd.ExcelWriter(output) as writer:
      pd.DataFrame({"info": ["No tables detected"]}).to_excel(writer, index=False, sheet_name="Tables")
    return
  with pd.ExcelWriter(output) as writer:
    for i, df in enumerate(dfs):
      sheet = f"Table_{i+1}"
      df.to_excel(writer, index=False, sheet_name=sheet)

def main():
  parser = argparse.ArgumentParser()
  parser.add_argument("op", type=str)
  parser.add_argument("--inputs", nargs="*")
  parser.add_argument("--images", nargs="*")
  parser.add_argument("--input")
  parser.add_argument("--output")
  parser.add_argument("--outdir")
  parser.add_argument("--ranges", default="")
  parser.add_argument("--rotation", type=int, default=90)
  parser.add_argument("--password", default="")
  parser.add_argument("--order", default="")
  parser.add_argument("--text", default="")
  parser.add_argument("--image", default="")
  args = parser.parse_args()

  if args.op == "merge":
    merge_pdfs(args.inputs or [], args.output)
  elif args.op == "split":
    split_pdf(args.input, args.ranges, args.outdir)
  elif args.op == "compress":
    compress_pdf(args.input, args.output)
  elif args.op == "pdf_to_word":
    pdf_to_word(args.input, args.output)
  elif args.op == "pdf_to_pptx":
    pdf_to_pptx(args.input, args.output)
  elif args.op == "jpg_to_pdf":
    jpg_to_pdf(args.images or [], args.output)
  elif args.op == "pdf_to_jpg":
    pdf_to_jpg(args.input, args.outdir)
  elif args.op == "rotate":
    rotate_pdf(args.input, args.output, args.rotation, args.ranges)
  elif args.op == "unlock":
    unlock_pdf(args.input, args.output, args.password)
  elif args.op == "protect":
    protect_pdf(args.input, args.output, args.password)
  elif args.op == "reorder_pages":
    order = []
    if args.order:
      try:
        order = json.loads(args.order)
      except:
        pass
    reorder_pages(args.input, args.output, order)
  elif args.op == "watermark_text":
    watermark_text(args.input, args.output, args.text)
  elif args.op == "watermark_image":
    watermark_image(args.input, args.output, args.image)
  elif args.op == "pdf_to_excel":
    pdf_to_excel(args.input, args.output)
  else:
    raise SystemExit(f"Unknown op: {args.op}")

if __name__ == "__main__":
  main()
